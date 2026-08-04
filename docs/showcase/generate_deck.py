"""Generate the Teams Agents technical showcase slide deck (.pptx).

Self-contained python-pptx generator. No external template needed.

Run:
    python docs/showcase/generate_deck.py

Output:
    docs/showcase/TeamsAgents-Technical-Showcase.pptx
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

# ----------------------------------------------------------------------------
# Palette
# ----------------------------------------------------------------------------
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
NAVY2 = RGBColor(0x12, 0x2C, 0x52)
PRIMARY = RGBColor(0x1F, 0x4E, 0x79)
OPT_A = RGBColor(0x2E, 0x75, 0xB6)  # blue = app-only
OPT_B = RGBColor(0x1F, 0x9E, 0x74)  # teal-green = per-user OBO
AMBER = RGBColor(0xE8, 0xA3, 0x3D)  # gotcha / HITL
RED = RGBColor(0xC0, 0x3B, 0x3B)
LIGHT = RGBColor(0xF2, 0xF5, 0xF9)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x6B, 0x72, 0x80)
DARKTEXT = RGBColor(0x1A, 0x20, 0x2C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ROW_ALT = RGBColor(0xED, 0xF1, 0xF7)

FONT = "Segoe UI"
FONT_LIGHT = "Segoe UI Light"
MONO = "Consolas"

# Slide geometry (16:9)
SW = Emu(12192000)
SH = Emu(6858000)


def _in(v: float) -> Emu:
    return Emu(int(v * 914400))


prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# ----------------------------------------------------------------------------
# Low-level helpers
# ----------------------------------------------------------------------------
def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def _set_line(shape, color, width_pt=1.0):
    ln = shape.line
    ln.color.rgb = color
    ln.width = Pt(width_pt)


def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE,
         radius=None):
    sp = s.shapes.add_shape(shape, _in(x), _in(y), _in(w), _in(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        _set_line(sp, line, line_w)
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def _apply_runs(tf, lines, default_size, color, align, font=FONT):
    """lines: list of dicts {text, size, bold, color, italic, bullet, level, space_after}."""
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", align)
        p.level = ln.get("level", 0)
        if ln.get("space_before") is not None:
            p.space_before = Pt(ln["space_before"])
        p.space_after = Pt(ln.get("space_after", 4))
        segs = ln.get("segs")
        if segs is None:
            segs = [{"text": ln.get("text", ""), "bold": ln.get("bold", False),
                     "color": ln.get("color", color), "size": ln.get("size", default_size),
                     "italic": ln.get("italic", False), "font": ln.get("font", font)}]
        for j, seg in enumerate(segs):
            r = p.add_run() if (j > 0 or p.runs) else p.add_run()
            r.text = seg["text"]
            r.font.name = seg.get("font", font)
            r.font.size = Pt(seg.get("size", ln.get("size", default_size)))
            r.font.bold = seg.get("bold", False)
            r.font.italic = seg.get("italic", False)
            r.font.color.rgb = seg.get("color", color)


def text(s, x, y, w, h, lines, size=16, color=DARKTEXT, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, font=FONT, wrap=True):
    tb = s.shapes.add_textbox(_in(x), _in(y), _in(w), _in(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if isinstance(lines, str):
        lines = [{"text": lines}]
    _apply_runs(tf, lines, size, color, align, font)
    return tb


def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt


# ----------------------------------------------------------------------------
# Slide furniture
# ----------------------------------------------------------------------------
def header(s, kicker, title, accent=PRIMARY):
    rect(s, 0, 0, 13.333, 1.28, fill=NAVY)
    rect(s, 0, 1.28, 13.333, 0.06, fill=accent)
    rect(s, 0.55, 0.30, 0.14, 0.66, fill=accent)
    text(s, 0.85, 0.24, 11.8, 0.34, kicker.upper(), size=12.5, color=accent,
         font=FONT)
    text(s, 0.85, 0.55, 11.8, 0.62, [{"text": title, "bold": True, "size": 26,
                                      "color": WHITE}], font=FONT)


def footer(s, idx):
    text(s, 0.55, 7.02, 8, 0.3,
         "One deployed agent \u00b7 two front doors \u00b7 Foundry hosted agent on Teams",
         size=9, color=GREY)
    text(s, 11.9, 7.02, 0.9, 0.3, str(idx), size=9, color=GREY, align=PP_ALIGN.RIGHT)


def chip(s, x, y, label, color, w=1.9, h=0.42, tcolor=WHITE, size=12):
    c = rect(s, x, y, w, h, fill=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    text(s, x, y, w, h, [{"text": label, "bold": True, "size": size, "color": tcolor}],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return c


def legend_dot(s, x, y, color, label):
    rect(s, x, y + 0.03, 0.22, 0.22, fill=color, shape=MSO_SHAPE.OVAL)
    text(s, x + 0.32, y, 3.2, 0.3, [{"text": label, "size": 11.5, "color": DARKTEXT}])


# ----------------------------------------------------------------------------
# Flow diagram: vertical numbered steps
# ----------------------------------------------------------------------------
def flow(s, x, y, w, steps, accent, step_h=0.62, gap=0.12, num_from=1):
    """steps: list of (actor, action) or (actor, action, color)."""
    cy = y
    for i, st in enumerate(steps):
        actor, action = st[0], st[1]
        col = st[2] if len(st) > 2 else accent
        # number badge
        rect(s, x, cy, 0.5, step_h, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.35)
        text(s, x, cy, 0.5, step_h, [{"text": str(num_from + i), "bold": True,
             "size": 15, "color": WHITE}], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
        # step card
        rect(s, x + 0.6, cy, w - 0.6, step_h, fill=CARD, line=RGBColor(0xD5, 0xDD, 0xE6),
             line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
        rect(s, x + 0.6, cy, 0.08, step_h, fill=col)
        text(s, x + 0.82, cy + 0.03, w - 0.95, step_h - 0.06,
             [{"segs": [{"text": actor + "  ", "bold": True, "size": 12.5, "color": col}]},
              {"segs": [{"text": action, "size": 11, "color": DARKTEXT}]}],
             anchor=MSO_ANCHOR.MIDDLE)
        cy += step_h + gap
    return cy


def card_box(s, x, y, w, h, title, body_lines, accent, title_size=14, body_size=11.5):
    rect(s, x, y, w, h, fill=CARD, line=RGBColor(0xD5, 0xDD, 0xE6), line_w=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    rect(s, x, y, w, 0.5, fill=accent, shape=MSO_SHAPE.RECTANGLE)
    text(s, x + 0.2, y, w - 0.4, 0.5, [{"text": title, "bold": True, "size": title_size,
         "color": WHITE}], anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.24, y + 0.62, w - 0.48, h - 0.7, body_lines, size=body_size,
         color=DARKTEXT, anchor=MSO_ANCHOR.TOP)


# ----------------------------------------------------------------------------
# Table
# ----------------------------------------------------------------------------
def table(s, x, y, w, h, headers, rows, col_ratios=None, header_fill=PRIMARY,
          font_size=11, header_size=11.5, first_col_bold=True, cell_colors=None):
    nrows = len(rows) + 1
    ncols = len(headers)
    gtbl = s.shapes.add_table(nrows, ncols, _in(x), _in(y), _in(w), _in(h)).table
    # disable banding style default
    gtbl.first_row = False
    gtbl.horz_banding = False
    # column widths
    if col_ratios is None:
        col_ratios = [1] * ncols
    tot = sum(col_ratios)
    for ci, r in enumerate(col_ratios):
        gtbl.columns[ci].width = _in(w * r / tot)
    # header
    for ci, htext in enumerate(headers):
        cell = gtbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = _in(0.08)
        cell.margin_right = _in(0.06)
        cell.margin_top = _in(0.03)
        cell.margin_bottom = _in(0.03)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
        r = p.add_run()
        r.text = htext
        r.font.name = FONT
        r.font.size = Pt(header_size)
        r.font.bold = True
        r.font.color.rgb = WHITE
    # body
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gtbl.cell(ri + 1, ci)
            cell.fill.solid()
            if cell_colors and cell_colors.get((ri, ci)) is not None:
                cell.fill.fore_color.rgb = cell_colors[(ri, ci)]
            else:
                cell.fill.fore_color.rgb = CARD if ri % 2 == 0 else ROW_ALT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = _in(0.08)
            cell.margin_right = _in(0.06)
            cell.margin_top = _in(0.02)
            cell.margin_bottom = _in(0.02)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = val
            run.font.name = FONT
            run.font.size = Pt(font_size)
            run.font.bold = bool(first_col_bold and ci == 0)
            run.font.color.rgb = DARKTEXT
    return gtbl


# ============================================================================
# SLIDE 1 — Title
# ============================================================================
def s_title():
    s = slide()
    bg(s, NAVY)
    rect(s, 0, 0, 13.333, 7.5, fill=NAVY)
    # accent bands
    rect(s, 0, 6.9, 13.333, 0.6, fill=NAVY2)
    rect(s, 0.0, 2.62, 0.9, 0.1, fill=OPT_A)
    text(s, 0.9, 0.9, 11.5, 0.4, "TECHNICAL SHOWCASE  \u00b7  CAPITAL MARKETS",
         size=15, color=OPT_B, font=FONT)
    text(s, 0.88, 1.7, 11.6, 1.7,
         [{"segs": [{"text": "Foundry Hosted Agents on Microsoft Teams",
                     "bold": True, "size": 40, "color": WHITE, "font": FONT}]},
          {"segs": [{"text": "Option A vs Option B \u2014 app-only vs per-user OBO",
                     "size": 24, "color": RGBColor(0xBF, 0xD3, 0xE8), "font": FONT_LIGHT}],
           "space_before": 8}])
    # front-door chips
    chip(s, 0.9, 3.85, "Microsoft Teams", OPT_A, w=2.5, h=0.5, size=13)
    chip(s, 3.6, 3.85, "Custom Web App", OPT_A, w=2.5, h=0.5, size=13)
    text(s, 6.35, 3.9, 6, 0.45, [{"text": "\u2192  one deployed agent, two front doors",
         "italic": True, "size": 15, "color": RGBColor(0xBF, 0xD3, 0xE8)}],
         anchor=MSO_ANCHOR.MIDDLE)
    # tagline
    text(s, 0.9, 5.0, 11.5, 1.4,
         [{"text": "Same downstream Azure AI Search call from both front doors \u2014 the only",
           "size": 15, "color": RGBColor(0x9F, 0xB4, 0xCC)},
          {"text": "difference is whose identity drives document-level trimming.",
           "size": 15, "color": RGBColor(0x9F, 0xB4, 0xCC)}])
    text(s, 0.9, 6.98, 11.5, 0.4,
         "FastAPI backend  \u00b7  React SPA  \u00b7  Foundry hosted agent  \u00b7  M365 Agents SDK proxy",
         size=11.5, color=GREY)
    notes(s, "Welcome. This deck walks a technical audience through two supported ways to "
              "expose a Foundry hosted agent in Teams: Option A (direct publish, app-only "
              "search) and Option B (Custom Engine Agent proxy doing Teams SSO + OBO for "
              "per-user document security). Both also work from a custom web app. The core "
              "design principle: one deployed agent, two front doors, converging on the same "
              "Azure AI Search call.")


# ============================================================================
# SLIDE 2 — Agenda
# ============================================================================
def s_agenda():
    s = slide()
    bg(s, LIGHT)
    header(s, "Agenda", "What we will cover")
    items = [
        ("01", "The big picture", "One agent, two front doors, one Search call", OPT_A),
        ("02", "Option A end-to-end", "App-only \u2014 Teams and Web flows", OPT_A),
        ("03", "Option B end-to-end", "Per-user OBO \u2014 Teams and Web flows", OPT_B),
        ("04", "Technical deep dive", "Who does the OBO, token audiences, native ACL", PRIMARY),
        ("05", "Options & feature matrix", "A vs B vs C vs D \u2014 pros / cons / benefits", AMBER),
        ("06", "Showcase runbook", "Live demo script + key takeaways", OPT_B),
    ]
    x0, y0 = 0.7, 1.75
    cw, ch = 5.9, 1.45
    gx, gy = 0.35, 0.28
    for i, (num, t, sub, col) in enumerate(items):
        cx = x0 + (i % 2) * (cw + gx)
        cy = y0 + (i // 2) * (ch + gy)
        rect(s, cx, cy, cw, ch, fill=CARD, line=RGBColor(0xDD, 0xE3, 0xEB), line_w=1.0,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
        rect(s, cx, cy, 0.12, ch, fill=col)
        text(s, cx + 0.32, cy + 0.18, 1.4, 1.1,
             [{"text": num, "bold": True, "size": 34, "color": col}],
             anchor=MSO_ANCHOR.MIDDLE)
        text(s, cx + 1.55, cy + 0.24, cw - 1.75, 0.5,
             [{"text": t, "bold": True, "size": 17, "color": DARKTEXT}])
        text(s, cx + 1.55, cy + 0.74, cw - 1.75, 0.6,
             [{"text": sub, "size": 12.5, "color": GREY}])
    footer(s, 2)
    notes(s, "Six parts: framing, the two options end to end, the technical mechanics, the "
              "decision matrix, and a live demo runbook with takeaways.")


# ============================================================================
# SLIDE 3 — Big picture
# ============================================================================
def s_bigpicture():
    s = slide()
    bg(s, LIGHT)
    header(s, "The big picture", "One deployed agent \u00b7 two front doors")
    # two front doors
    card_box(s, 0.7, 1.75, 3.4, 1.9, "Microsoft Teams", [
        {"text": "\u2022 Native Teams chat client", "size": 12},
        {"text": "\u2022 Option A: direct Bot Service publish", "size": 12},
        {"text": "\u2022 Option B: CEA proxy (SSO + OBO)", "size": 12},
    ], OPT_A)
    card_box(s, 0.7, 3.9, 3.4, 1.9, "Custom Web App (SPA)", [
        {"text": "\u2022 React + MSAL.js", "size": 12},
        {"text": "\u2022 Option A: app-only endpoint", "size": 12},
        {"text": "\u2022 Option B: user token \u2192 backend OBO", "size": 12},
    ], OPT_A)
    # convergence
    rect(s, 4.5, 2.6, 0.7, 0.02, fill=GREY)
    rect(s, 4.5, 4.9, 0.7, 0.02, fill=GREY)
    card_box(s, 5.25, 2.55, 3.5, 2.5, "auth_context.from_headers()", [
        {"text": "Both front doors converge to a single", "size": 12},
        {"text": "UserAuth object:", "size": 12},
        {"text": "\u2022 search_token  (already exchanged)", "size": 11.5,
         "color": OPT_B},
        {"text": "\u2022 user_assertion \u2192 backend OBO", "size": 11.5, "color": OPT_B},
        {"text": " ", "size": 6},
        {"text": "Option A skips identity entirely \u2192", "size": 11.5, "color": OPT_A},
        {"text": "application / admin identity.", "size": 11.5, "color": OPT_A},
    ], PRIMARY)
    # search
    card_box(s, 9.15, 2.9, 3.45, 1.8, "Azure AI Search", [
        {"text": "Document-level security", "size": 12, "bold": True},
        {"text": "\u2022 Native ACL trims to caller's", "size": 11.5},
        {"text": "   real Entra groups / oid", "size": 11.5},
        {"text": "\u2022 One index, one query shape", "size": 11.5},
    ], OPT_B)
    rect(s, 8.75, 3.75, 0.4, 0.02, fill=GREY)
    text(s, 0.7, 6.05, 12, 0.8,
         [{"segs": [{"text": "Key idea:  ", "bold": True, "size": 14, "color": PRIMARY},
                    {"text": "whether a request arrives from Teams or the web app, it resolves "
                     "to the same downstream AI Search call. The only variable is ",
                     "size": 14, "color": DARKTEXT},
                    {"text": "whose identity", "bold": True, "italic": True, "size": 14,
                     "color": OPT_B},
                    {"text": " drives trimming.", "size": 14, "color": DARKTEXT}]}])
    footer(s, 3)
    notes(s, "The mental model for the whole deck. Two front doors (Teams, Web), each with an "
              "Option A and Option B path. Option B paths both converge on auth_context -> "
              "UserAuth -> per-user Search. Option A skips identity and uses the app/admin "
              "identity. This is why a single deployed agent can serve both surfaces.")


# ============================================================================
# SLIDE 4 — Option A overview
# ============================================================================
def s_optA_overview():
    s = slide()
    bg(s, LIGHT)
    header(s, "Option A", "Direct Foundry publish \u2014 app-only", accent=OPT_A)
    text(s, 0.7, 1.6, 11.9, 0.7,
         [{"segs": [{"text": "The hosted agent is published straight to Teams via Azure Bot "
                     "Service. Its in-container Search tool runs with the ", "size": 15,
                     "color": DARKTEXT},
                    {"text": "application identity", "bold": True, "size": 15, "color": OPT_A},
                    {"text": " \u2014 every user sees the same result set. No per-user trimming.",
                     "size": 15, "color": DARKTEXT}]}])
    card_box(s, 0.7, 2.5, 3.85, 3.9, "How it works", [
        {"text": "\u2022 Foundry portal / REST \u2192 Publish to", "size": 12.5},
        {"text": "   Teams & M365 Copilot", "size": 12.5},
        {"text": "\u2022 Azure Bot Service bridges", "size": 12.5},
        {"text": "   Responses \u2194 Activity automatically", "size": 12.5},
        {"text": "\u2022 In-container tool uses", "size": 12.5},
        {"text": "   DefaultAzureCredential = agent MI", "size": 12.5},
        {"text": " ", "size": 6},
        {"text": "Availability: GA publish flow", "size": 12.5, "bold": True, "color": OPT_B},
    ], OPT_A)
    card_box(s, 4.75, 2.5, 3.85, 3.9, "Strengths", [
        {"text": "\u2713 Lowest effort \u2014 portal / REST", "size": 12.5, "color": OPT_B},
        {"text": "\u2713 Foundry-native tracing + eval", "size": 12.5, "color": OPT_B},
        {"text": "\u2713 Dedicated Entra agent identity", "size": 12.5, "color": OPT_B},
        {"text": "\u2713 Stable endpoint + version selector", "size": 12.5, "color": OPT_B},
        {"text": "\u2713 Inline citations render natively", "size": 12.5, "color": OPT_B},
        {"text": "\u2713 Tenant scope \u2192 M365 admin approval", "size": 12.5, "color": OPT_B},
    ], PRIMARY)
    card_box(s, 8.75, 2.5, 3.85, 3.9, "Limits", [
        {"text": "\u2717 No per-user document trimming", "size": 12.5, "color": RED},
        {"text": "\u2717 No OBO to downstream as user", "size": 12.5, "color": RED},
        {"text": "\u2717 No Adaptive Cards / buttons", "size": 12.5, "color": RED},
        {"text": "\u2717 No streaming, no card citations", "size": 12.5, "color": RED},
        {"text": " ", "size": 6},
        {"text": "Best for: internal Q&A / research", "size": 12.5, "bold": True,
         "color": AMBER},
        {"text": "where all users share entitlements.", "size": 12.5, "color": AMBER},
    ], GREY)
    footer(s, 4)
    notes(s, "Option A is the recommended default for the bulk of agents: GA, minimal glue, "
              "Foundry-native observability. The catch is app-only search - no per-user "
              "trimming and no user-context OBO - because the hosting gateway strips the "
              "inbound user token before it reaches the container.")


# ============================================================================
# SLIDE 5 — Option A Teams flow
# ============================================================================
def s_optA_teams():
    s = slide()
    bg(s, LIGHT)
    header(s, "Option A \u00b7 End-to-end", "Teams client \u2014 app-only", accent=OPT_A)
    steps = [
        ("Teams user", "\u201cWhat is our semiconductor view?\u201d \u2192 Activity message"),
        ("Azure Bot Service", "Bridges Responses \u2194 Activity (automatic, text only)"),
        ("Foundry hosted agent", "capmarkets-research-agent runs the Search tool"),
        ("Azure AI Search", "App-only query with the container's managed identity",),
        ("Hosted agent", "Grounds answer on shared results \u2192 text reply"),
        ("Teams user", "Same answer for everyone \u2014 no per-user trimming", AMBER),
    ]
    flow(s, 0.7, 1.7, 7.4, steps, OPT_A, step_h=0.7, gap=0.14)
    # side panel
    card_box(s, 8.5, 1.7, 4.1, 4.55, "Identity basis", [
        {"text": "Application / admin", "size": 13, "bold": True, "color": OPT_A},
        {"text": "Every user gets the same documents.", "size": 11.5},
        {"text": " ", "size": 6},
        {"text": "Watch-outs for the demo", "size": 12.5, "bold": True, "color": AMBER},
        {"text": "\u2022 The \u201cFoundry login\u201d prompt is NOT", "size": 11},
        {"text": "   an OBO \u2014 it just runs the agent.", "size": 11},
        {"text": "\u2022 Index has native ACL enabled, so a", "size": 11},
        {"text": "   token-less query matches 0 docs \u2014", "size": 11},
        {"text": "   even the public one.", "size": 11},
        {"text": "\u2022 Expect: \u201cno entitled research.\u201d", "size": 11,
         "italic": True, "color": RED},
        {"text": "   That gap is the teaching point.", "size": 11, "italic": True},
    ], PRIMARY)
    footer(s, 5)
    notes(s, "Teams Option A: user -> Bot Service -> hosted agent -> app-only Search. The "
              "Responses-to-Activity bridge preserves text but not citations, streaming, or "
              "cards. Because the index has native ACL enabled and no user token flows to the "
              "container, the published agent legitimately returns zero documents and answers "
              "'no entitled research' - which is exactly the gap Option B closes.")


# ============================================================================
# SLIDE 6 — Option A Web flow
# ============================================================================
def s_optA_web():
    s = slide()
    bg(s, LIGHT)
    header(s, "Option A \u00b7 End-to-end", "Web app \u2014 app-only", accent=OPT_A)
    steps = [
        ("Web user", "Select persona, ask a research query (no user token sent)"),
        ("Web SPA", "POST optionA/invoke"),
        ("Backend", "app_only_search \u2014 admin token OR classification ne 'mnpi'"),
        ("Azure AI Search", "Returns a baseline / non-MNPI slice"),
        ("Foundry hosted agent", "synthesize(answer, hits)"),
        ("Web user", "Answer \u2014 persona selector is illustrative only", AMBER),
    ]
    flow(s, 0.7, 1.7, 7.4, steps, OPT_A, step_h=0.7, gap=0.14)
    card_box(s, 8.5, 1.7, 4.1, 4.55, "Same label, different retrieval", [
        {"text": "Web Option A \u2260 Teams Option A", "size": 12.5, "bold": True,
         "color": AMBER},
        {"text": " ", "size": 5},
        {"text": "Teams (published agent):", "size": 11.5, "bold": True},
        {"text": "container MI \u2192 native ACL \u2192 0 docs", "size": 11},
        {"text": " ", "size": 4},
        {"text": "Web (optionA/invoke):", "size": 11.5, "bold": True},
        {"text": "admin token, OR a plain", "size": 11},
        {"text": "classification ne 'mnpi' filter \u2192", "size": 11, "font": MONO},
        {"text": "baseline / non-MNPI docs", "size": 11},
        {"text": " ", "size": 6},
        {"text": "Do not expect the two to match.", "size": 11, "italic": True,
         "color": RED},
    ], PRIMARY)
    footer(s, 6)
    notes(s, "Web Option A does NOT go token-less into native ACL. It either presents the "
              "backend admin managed-identity token (full entitlements) or applies a plain "
              "classification-ne-mnpi filter. So the web-app Option A returns a baseline set, "
              "while the Teams-published Option A returns zero. Same 'Option A' label, two "
              "different retrieval implementations.")


# ============================================================================
# SLIDE 7 — Option A technical details
# ============================================================================
def s_optA_tech():
    s = slide()
    bg(s, LIGHT)
    header(s, "Option A \u00b7 Technical details", "Why per-user OBO is structurally impossible",
           accent=OPT_A)
    card_box(s, 0.7, 1.65, 5.85, 2.35, "1 \u00b7 Gateway strips the user token", [
        {"text": "The hosted-agent runtime terminates TLS and", "size": 12},
        {"text": "injects only FOUNDRY_* env vars. The caller's", "size": 12},
        {"text": "Authorization bearer never reaches the", "size": 12, "font": MONO},
        {"text": "container, so the tool has no assertion to", "size": 12},
        {"text": "exchange \u2192 DefaultAzureCredential() =", "size": 12, "font": MONO},
        {"text": "the agent's managed identity.", "size": 12},
    ], OPT_A)
    card_box(s, 6.75, 1.65, 5.85, 2.35, "2 \u00b7 Bot Service yields no user token", [
        {"text": "Direct Teams publish authenticates via a", "size": 12},
        {"text": "channel JWT \u2014 a channel/authorization", "size": 12},
        {"text": "token, not an Entra user token. Only the", "size": 12},
        {"text": "Entra auth scheme (the proxy path) derives", "size": 12},
        {"text": "an OBO-exchangeable user identity.", "size": 12},
        {"text": "x-ms-user-identity is an isolation key only.", "size": 11, "font": MONO,
         "color": GREY},
    ], PRIMARY)
    # chain strip
    rect(s, 0.7, 4.35, 11.9, 1.05, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, 0.95, 4.4, 11.4, 0.95,
         [{"segs": [{"text": "Chain:  ", "bold": True, "size": 13, "color": AMBER},
                    {"text": "no OBO on Option A  \u2192  Search runs as the container "
                     "identity  \u2192  native ACL finds no matching membership  \u2192  "
                     "0 docs  \u2192  \u201cno entitled research.\u201d", "size": 13,
                     "color": WHITE}]}],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, 0.7, 5.65, 11.9, 0.9,
         [{"segs": [{"text": "The lesson:  ", "bold": True, "size": 14, "color": OPT_A},
                    {"text": "in a document-security-enabled index, an agent that can't flow "
                     "the user's identity to the tool returns nothing per user. That is exactly "
                     "the gap ", "size": 14, "color": DARKTEXT},
                    {"text": "Option B", "bold": True, "size": 14, "color": OPT_B},
                    {"text": " (proxy SSO + OBO) closes.", "size": 14, "color": DARKTEXT}]}])
    footer(s, 7)
    notes(s, "Two independent reasons per-user OBO can't work inside a directly-published "
              "hosted agent: (1) the runtime gateway strips the inbound bearer, and (2) the "
              "Bot Service channel JWT is not an Entra user token. The x-ms-user-identity "
              "header is an opaque session-isolation key, not an OBO token. Result chain shown "
              "in the dark strip.")


# ============================================================================
# SLIDE 8 — Option B overview
# ============================================================================
def s_optB_overview():
    s = slide()
    bg(s, LIGHT)
    header(s, "Option B", "Custom Engine Agent proxy \u2014 per-user OBO", accent=OPT_B)
    text(s, 0.7, 1.6, 11.9, 0.7,
         [{"segs": [{"text": "A proxy (M365 Agents SDK) is the ", "size": 15,
                     "color": DARKTEXT},
                    {"text": "trust boundary", "bold": True, "size": 15, "color": OPT_B},
                    {"text": ": Teams SSO \u2192 OBO \u2192 a per-user Azure AI Search token, so "
                     "results are trimmed to that user\u2019s entitlements (document-level "
                     "security).", "size": 15, "color": DARKTEXT}]}])
    card_box(s, 0.7, 2.5, 3.85, 3.9, "How it works", [
        {"text": "\u2022 Teams SSO signs the user in", "size": 12.5},
        {"text": "\u2022 Proxy exchanges the token OBO for", "size": 12.5},
        {"text": "   a search.azure.com token", "size": 12.5, "font": MONO},
        {"text": "\u2022 Forwards it to the backend via", "size": 12.5},
        {"text": "   x-ms-query-source-authorization", "size": 11, "font": MONO},
        {"text": "\u2022 Backend runs per-user Search", "size": 12.5},
        {"text": " ", "size": 6},
        {"text": "Availability: GA tooling", "size": 12.5, "bold": True, "color": OPT_B},
    ], OPT_B)
    card_box(s, 4.75, 2.5, 3.85, 3.9, "Strengths", [
        {"text": "\u2713 Real per-user doc trimming", "size": 12.5, "color": OPT_B},
        {"text": "\u2713 OBO to downstream as the user", "size": 12.5, "color": OPT_B},
        {"text": "\u2713 Adaptive Cards, buttons, HITL", "size": 12.5, "color": OPT_B},
        {"text": "\u2713 Full control of Teams UX", "size": 12.5, "color": OPT_B},
        {"text": "\u2713 Multi-channel (Teams, Web, +)", "size": 12.5, "color": OPT_B},
        {"text": "\u2713 Built-in SSO + CI/CD templates", "size": 12.5, "color": OPT_B},
    ], PRIMARY)
    card_box(s, 8.75, 2.5, 3.85, 3.9, "Costs", [
        {"text": "\u2717 Build + host a proxy app", "size": 12.5, "color": RED},
        {"text": "\u2717 Two hops to trace (bot + agent)", "size": 12.5, "color": RED},
        {"text": "\u2717 No GCC / Government publish", "size": 12.5, "color": RED},
        {"text": "\u2717 JS/TS today (Python planned)", "size": 12.5, "color": RED},
        {"text": " ", "size": 6},
        {"text": "Best for: regulated data, per-user", "size": 12.5, "bold": True,
         "color": AMBER},
        {"text": "security, rich UX, custom SSO/OBO.", "size": 12.5, "color": AMBER},
    ], GREY)
    footer(s, 8)
    notes(s, "Option B is the pattern for real document-level security and rich Teams UX. The "
              "proxy is the trust boundary that runs SSO and the OBO exchange. Cost is that you "
              "own and host an app, tracing spans two hops, and Agents Toolkit publishing isn't "
              "supported in Government tenants.")


# ============================================================================
# SLIDE 9 — Option B Teams flow
# ============================================================================
def s_optB_teams():
    s = slide()
    bg(s, LIGHT)
    header(s, "Option B \u00b7 End-to-end", "Teams client \u2014 per-user OBO", accent=OPT_B)
    steps = [
        ("Teams user", "Asks a research question \u2192 Activity to the CEA proxy"),
        ("Teams SSO", "Silent sign-in, aud = api://botid-{clientId}"),
        ("Bot Service", "getToken \u2192 user token (app's own identity)"),
        ("CEA proxy", "exchangeToken OBO \u2192 Azure AI Search token (this user)"),
        ("Backend optionB/invoke", "x-ms-query-source-authorization: <search token>"),
        ("Azure AI Search", "Per-user query \u2192 entitled documents only"),
        ("Teams user", "Adaptive Card: citations + \u201ctrimmed for you\u201d", OPT_B),
    ]
    flow(s, 0.7, 1.62, 7.55, steps, OPT_B, step_h=0.6, gap=0.11)
    card_box(s, 8.6, 1.62, 4.0, 4.75, "The proxy does all token work", [
        {"text": "Backend performs NO OBO here \u2014 it", "size": 11.5},
        {"text": "forwards the already-exchanged token.", "size": 11.5},
        {"text": " ", "size": 5},
        {"text": "Audience progression", "size": 12, "bold": True, "color": OPT_B},
        {"text": "SSO  \u2192 api://botid-{clientId}", "size": 10.5, "font": MONO},
        {"text": "OBO  \u2192 https://search.azure.com", "size": 10.5, "font": MONO},
        {"text": "\u2192 AI Search native ACL trims", "size": 10.5, "font": MONO},
        {"text": " ", "size": 5},
        {"text": "Fails closed", "size": 12, "bold": True, "color": AMBER},
        {"text": "No token \u2192 public-only + \u201csign in\u201d", "size": 11},
    ], PRIMARY)
    footer(s, 9)
    notes(s, "Teams Option B: the proxy is the trust boundary. It runs silent Teams SSO "
              "(audience api://botid-clientId), then exchangeToken performs the OBO to a "
              "search.azure.com-scoped token for this user, and POSTs it to the backend in "
              "x-ms-query-source-authorization. The backend does no OBO on this path - it uses "
              "the token directly. If SSO/OBO isn't complete, it fails closed to public-only.")


# ============================================================================
# SLIDE 10 — Option B Web flow
# ============================================================================
def s_optB_web():
    s = slide()
    bg(s, LIGHT)
    header(s, "Option B \u00b7 End-to-end", "Web app \u2014 per-user OBO", accent=OPT_B)
    steps = [
        ("Web user", "Sign in with MSAL.js"),
        ("Microsoft Entra ID", "acquireTokenSilent \u2192 user access token (backend API scope)"),
        ("Web SPA", "POST optionB/invoke \u2014 Authorization: Bearer <user token>"),
        ("Backend", "OBO exchange: user assertion \u2192 search.azure.com token"),
        ("Azure AI Search", "Per-user query (user's Entra groups / oid)"),
        ("Web user", "Per-user answer + entitled docs only", OPT_B),
    ]
    flow(s, 0.7, 1.7, 7.55, steps, OPT_B, step_h=0.7, gap=0.14)
    card_box(s, 8.6, 1.7, 4.0, 4.55, "Where OBO happens differs", [
        {"text": "Web path:", "size": 12, "bold": True, "color": OPT_B},
        {"text": "the BACKEND does the OBO", "size": 11.5},
        {"text": "(obo_service). SPA only ever sends", "size": 11, "font": MONO},
        {"text": "its own user token.", "size": 11},
        {"text": " ", "size": 6},
        {"text": "The browser never holds OBO", "size": 11.5, "bold": True, "color": AMBER},
        {"text": "credentials. auth_context puts the", "size": 11},
        {"text": "Bearer into user_assertion, which", "size": 11, "font": MONO},
        {"text": "triggers the OBO exchange.", "size": 11},
        {"text": " ", "size": 6},
        {"text": "Same per_user_search call as Teams.", "size": 11.5, "italic": True,
         "color": PRIMARY},
    ], PRIMARY)
    footer(s, 10)
    notes(s, "Web Option B: the SPA signs the user in with MSAL and sends the user's own token "
              "as a Bearer header for the backend API scope. The BACKEND (not the browser) "
              "performs the On-Behalf-Of exchange, then runs per-user Search. Same trimming "
              "semantics and same per_user_search call as Teams Option B.")


# ============================================================================
# SLIDE 11 — Convergence: who does the OBO
# ============================================================================
def s_converge():
    s = slide()
    bg(s, LIGHT)
    header(s, "Option B \u00b7 The one contrast to memorize",
           "Where the OBO happens: Teams vs Web")
    table(s, 0.7, 1.7, 11.9, 1.9,
          ["Path", "Who performs the OBO", "What the backend receives", "Header used"],
          [["Teams (Option B)", "The proxy (exchangeToken)", "An already-exchanged Search token",
            "x-ms-query-source-authorization"],
           ["Web app (Option B)", "The backend (obo_service)", "The user's own token (needs OBO)",
            "Authorization: Bearer"]],
          col_ratios=[1.3, 1.6, 2.0, 1.9], header_fill=PRIMARY, font_size=12.5)
    # convergence diagram
    chip(s, 0.9, 4.15, "CEA proxy \u2014 SSO + OBO", OPT_B, w=3.1, h=0.55, size=12)
    chip(s, 0.9, 5.0, "SPA MSAL \u2014 Bearer token", OPT_A, w=3.1, h=0.55, size=12)
    rect(s, 4.05, 4.42, 1.15, 0.02, fill=GREY)
    rect(s, 4.05, 5.27, 1.15, 0.02, fill=GREY)
    card_box(s, 5.25, 4.05, 3.2, 1.75, "auth_context\n.from_headers()", [
        {"text": "\u2192 single UserAuth", "size": 12, "bold": True},
        {"text": "search_token", "size": 11, "font": MONO, "color": OPT_B},
        {"text": "user_assertion \u2192 OBO", "size": 11, "font": MONO, "color": OPT_B},
    ], PRIMARY, title_size=12)
    rect(s, 8.45, 4.85, 0.8, 0.02, fill=GREY)
    card_box(s, 9.3, 4.05, 3.3, 1.75, "per_user_search", [
        {"text": "Identical call for both paths.", "size": 11.5},
        {"text": "Azure AI Search enforces", "size": 11.5},
        {"text": "document-level security.", "size": 11.5, "bold": True, "color": OPT_B},
    ], OPT_B, title_size=12)
    text(s, 0.7, 6.35, 11.9, 0.6,
         [{"segs": [{"text": "One agent, two front doors:  ", "bold": True, "size": 13.5,
                     "color": PRIMARY},
                    {"text": "both Option B paths resolve to the identical per-user Search "
                     "call. The browser never holds OBO credentials.", "size": 13.5,
                     "color": DARKTEXT}]}])
    footer(s, 11)
    notes(s, "This is the single contrast worth memorizing. On Teams the PROXY does the OBO and "
              "the backend receives an already-exchanged Search token in "
              "x-ms-query-source-authorization. On Web the BACKEND does the OBO from the user's "
              "own Bearer token. Both converge on the identical per_user_search call.")


# ============================================================================
# SLIDE 12 — Option B technical details
# ============================================================================
def s_optB_tech():
    s = slide()
    bg(s, LIGHT)
    header(s, "Option B \u00b7 Technical details",
           "SSO \u2192 OBO token mechanics (verified live)", accent=OPT_B)
    card_box(s, 0.7, 1.65, 5.85, 4.6, "The four moving parts", [
        {"text": "1  SSO connection scope", "size": 12.5, "bold": True, "color": OPT_B},
        {"text": "search-sso scoped to the app's OWN", "size": 11, "font": MONO},
        {"text": "identity (.../access_as_user) so the base", "size": 11, "font": MONO},
        {"text": "token is OBO-exchangeable.", "size": 11},
        {"text": "2  tokenExchangeUrl = api://botid-{id}", "size": 12.5, "bold": True,
         "color": OPT_B},
        {"text": "Standalone bot SSO token aud carries the", "size": 11},
        {"text": "botid- prefix \u2014 must string-match manifest.", "size": 11},
        {"text": "3  exchangeToken \u2192 oboScopes", "size": 12.5, "bold": True,
         "color": OPT_B},
        {"text": "[search.azure.com/.default] performs the", "size": 11, "font": MONO},
        {"text": "OBO to a per-user Search token.", "size": 11},
        {"text": "4  Delegated consent granted", "size": 12.5, "bold": True, "color": OPT_B},
        {"text": "App has Azure Cognitive Search", "size": 11},
        {"text": "user_impersonation + admin consent.", "size": 11, "font": MONO},
    ], PRIMARY)
    card_box(s, 6.75, 1.65, 5.85, 2.2, "Token audience progression", [
        {"text": "Teams user session", "size": 12, "font": MONO},
        {"text": "  \u2192 (SSO)   aud = api://botid-{clientId}", "size": 11.5, "font": MONO,
         "color": OPT_A},
        {"text": "  \u2192 (OBO)   aud = https://search.azure.com", "size": 11.5,
         "font": MONO, "color": OPT_B},
        {"text": "  \u2192 AI Search \u2014 native ACL trims", "size": 11.5, "font": MONO},
    ], NAVY, title_size=13)
    card_box(s, 6.75, 4.05, 5.85, 2.2, "Native ACL vs GA trimming", [
        {"text": "Native ACL (use_native_acl=true):", "size": 11.5, "bold": True},
        {"text": "trims by the REAL signed-in user's oid /", "size": 11},
        {"text": "groups \u2014 persona dropdown has no effect.", "size": 11},
        {"text": "GA trimming (use_native_acl=false):", "size": 11.5, "bold": True},
        {"text": "builds filter from persona.entra_group_id \u2192", "size": 11, "font": MONO},
        {"text": "per-persona subsets for the same tester.", "size": 11},
    ], AMBER, title_size=13)
    footer(s, 12)
    notes(s, "The verified mechanics: the SSO connection must be scoped to the app's own "
              "identity (not search.azure.com) so its token is OBO-exchangeable; the "
              "tokenExchangeUrl and manifest resource must be the api://botid-{clientId} form "
              "for a standalone bot; exchangeToken with oboScopes=search.azure.com does the "
              "OBO; and the app needs delegated Search user_impersonation with admin consent. "
              "Native ACL trims by the real user; GA trimming trims by the selected persona.")


# ============================================================================
# SLIDE 13 — Options comparison matrix
# ============================================================================
def s_matrix():
    s = slide()
    bg(s, LIGHT)
    header(s, "Decision matrix", "Options A / B / C / D \u2014 pros, cons, benefits")
    headers = ["Aspect", "A \u00b7 Direct publish", "B \u00b7 Agents SDK proxy",
               "C \u00b7 Copilot Studio", "D \u00b7 A2A / MCP"]
    rows = [
        ["Availability", "GA", "GA tooling", "GA", "Preview"],
        ["Per-user doc trimming (OBO)", "\u2717", "\u2713 native", "\u25d1 partial", "\u25d1"],
        ["Teams UX (cards, buttons)", "\u2717 text only", "\u2713 full", "\u25d1", "\u25d1"],
        ["Dev effort", "\u2713 lowest", "\u2717 highest", "\u2713 low-code", "\u25d1"],
        ["Foundry-native tracing/eval", "\u2713", "\u25d1 two hops", "\u25d1", "\u25d1"],
        ["Identity / SSO / OBO", "\u25d1 channel auth", "\u2713 SSO + OBO", "\u25d1", "\u25d1"],
        ["Multi-channel beyond Teams", "\u25d1 Teams + M365", "\u2713 many", "\u25d1", "\u25d1"],
        ["Government / GCC publish", "\u2713", "\u2717 unsupported", "\u25d1 verify", "\u25d1"],
        ["Time to demo", "\u2713 fastest", "\u2717 slowest", "\u2713 fast", "\u25d1 medium"],
    ]
    # color the verdict cells lightly
    cc = {}
    green = RGBColor(0xDE, 0xF2, 0xE6)
    redc = RGBColor(0xF7, 0xDE, 0xDE)
    amb = RGBColor(0xFB, 0xEE, 0xD5)
    verdict = {
        (0, 1): green, (0, 2): green, (0, 3): green, (0, 4): amb,
        (1, 1): redc, (1, 2): green, (1, 3): amb, (1, 4): amb,
        (2, 1): redc, (2, 2): green, (2, 3): amb, (2, 4): amb,
        (3, 1): green, (3, 2): redc, (3, 3): green, (3, 4): amb,
        (4, 1): green, (4, 2): amb, (4, 3): amb, (4, 4): amb,
        (5, 1): amb, (5, 2): green, (5, 3): amb, (5, 4): amb,
        (6, 1): amb, (6, 2): green, (6, 3): amb, (6, 4): amb,
        (7, 1): green, (7, 2): redc, (7, 3): amb, (7, 4): amb,
        (8, 1): green, (8, 2): redc, (8, 3): green, (8, 4): amb,
    }
    table(s, 0.55, 1.62, 12.25, 4.55, headers, rows,
          col_ratios=[2.0, 1.55, 1.6, 1.5, 1.35], header_fill=PRIMARY,
          font_size=10.5, header_size=11, cell_colors=verdict)
    # legend
    legend_dot(s, 0.6, 6.35, RGBColor(0x2E, 0x8B, 0x57), "\u2713 strong / native")
    legend_dot(s, 3.4, 6.35, AMBER, "\u25d1 partial / needs work")
    legend_dot(s, 6.6, 6.35, RED, "\u2717 not supported / heavy lift")
    footer(s, 13)
    notes(s, "The full decision matrix. A wins on availability, effort, Foundry-native "
              "tracing, governance. B wins on per-user OBO, UX richness, SSO/CI-CD, "
              "multi-channel - at the cost of building/hosting an app and no GCC publish. C is "
              "low-code managed SaaS. D is a preview protocol surface, not a Teams client "
              "experience. Recommendation: default to A, reserve B for the subset needing "
              "per-user security or rich UX.")


# ============================================================================
# SLIDE 14 — Teams UX capability
# ============================================================================
def s_ux():
    s = slide()
    bg(s, LIGHT)
    header(s, "Teams UX & identity", "What Option A gives vs what needs Option B")
    table(s, 0.7, 1.7, 6.9, 3.9,
          ["Teams capability", "A", "B"],
          [["Conversational chat + Markdown", "\u2713", "\u2013"],
           ["Inline footnote citations", "\u2713", "\u2013"],
           ["Authored Adaptive Cards", "\u2717", "\u2713"],
           ["Buttons / Action.Execute (HITL)", "\u2717", "\u2713"],
           ["Clarification cards, suggested prompts", "\u2717", "\u2713"],
           ["Thumbs up/down feedback", "\u2717", "\u2713"],
           ["Message extensions, tabs, dialogs", "\u2717", "\u2713"],
           ["Per-user OBO to downstream", "\u2717", "\u2713"]],
          col_ratios=[3.4, 0.6, 0.6], header_fill=PRIMARY, font_size=11.5)
    card_box(s, 7.85, 1.7, 4.75, 3.9, "The nuance", [
        {"text": "Option A keeps conversational chat", "size": 12.5},
        {"text": "and native citations \u2014 but loses the", "size": 12.5},
        {"text": "authored UI and customer-controlled", "size": 12.5},
        {"text": "identity flows.", "size": 12.5},
        {"text": " ", "size": 6},
        {"text": "Direct publish provides channel-level", "size": 12},
        {"text": "auth + a dedicated agent identity for", "size": 12},
        {"text": "app-only downstream access.", "size": 12},
        {"text": " ", "size": 6},
        {"text": "Per-user OBO (attended flow) needs a", "size": 12, "bold": True,
         "color": OPT_B},
        {"text": "trust-boundary app \u2192 Option B.", "size": 12, "bold": True,
         "color": OPT_B},
    ], PRIMARY)
    text(s, 0.7, 5.85, 11.9, 0.9,
         [{"segs": [{"text": "Hybrid is common:  ", "bold": True, "size": 13.5,
                     "color": AMBER},
                    {"text": "default LOB agents on Option A; high-touch agents (trade approval "
                     "needing per-user Graph/OBO + rich cards) on Option B.", "size": 13.5,
                     "color": DARKTEXT}]}])
    footer(s, 14)
    notes(s, "Direct publish keeps conversational chat and native inline citations, but not "
              "authored Adaptive Cards, buttons, clarification cards, feedback controls, or "
              "message extensions/tabs - and not per-user OBO. Those require the Option B "
              "proxy. A hybrid split is the common recommendation.")


# ============================================================================
# SLIDE 15 — Security & identity deep dive
# ============================================================================
def s_security():
    s = slide()
    bg(s, LIGHT)
    header(s, "Security & identity", "Document-level trimming, verified behavior")
    card_box(s, 0.7, 1.65, 5.85, 2.5, "Native ACL (token-driven)", [
        {"text": "permissionFilterOption = enabled", "size": 11.5, "font": MONO},
        {"text": "Every query enforces GroupIds/UserIds.", "size": 12},
        {"text": "No user token \u2192 0 docs (even public).", "size": 12, "color": RED},
        {"text": "Real user token in", "size": 12},
        {"text": "x-ms-query-source-authorization \u2192 only", "size": 11, "font": MONO},
        {"text": "docs matching the caller's real Entra", "size": 12},
        {"text": "groups / oid. True per-user security.", "size": 12, "bold": True,
         "color": OPT_B},
    ], OPT_B)
    card_box(s, 6.75, 1.65, 5.85, 2.5, "GA security trimming (persona)", [
        {"text": "use_native_acl = false", "size": 11.5, "font": MONO},
        {"text": "Filter built from the persona's", "size": 12},
        {"text": "entra_group_id:", "size": 11.5, "font": MONO},
        {"text": "group_ids/any(g: search.in(g, ...))", "size": 10.5, "font": MONO},
        {"text": "Different personas \u2192 different subsets", "size": 12},
        {"text": "for the SAME tester. Empty groups fall", "size": 12},
        {"text": "back to public (GRP_ALL) only.", "size": 12},
    ], AMBER)
    card_box(s, 0.7, 4.35, 11.9, 2.05, "Verified live in this environment", [
        {"segs": [{"text": "\u2022  ", "size": 12, "color": OPT_B},
                  {"text": "Compliance user (oid 91a5dd7e = full control) sees all 8 docs "
                   "including MNPI (DEAL-MEMO-007, SURV-021).", "size": 12}]},
        {"segs": [{"text": "\u2022  ", "size": 12, "color": AMBER},
                  {"text": "Native-ACL subtlety: it trims by the real signed-in user, NOT the "
                   "selected persona \u2014 the tester always sees their full entitlement "
                   "regardless of dropdown.", "size": 12}]},
        {"segs": [{"text": "\u2022  ", "size": 12, "color": PRIMARY},
                  {"text": "To demo per-persona differences as one tester, use GA trimming; to "
                   "demo true token-driven security, sign in as separate group members.",
                   "size": 12}]},
        {"segs": [{"text": "\u2022  ", "size": 12, "color": RED},
                  {"text": "Docs must carry REAL tenant group object-ids / user oids \u2014 "
                   "synthetic GUIDs match nothing under native ACL.", "size": 12}]},
    ], NAVY, title_size=14)
    footer(s, 15)
    notes(s, "Two trimming modes. Native ACL is token-driven: no user token means zero docs, a "
              "real user token means only the caller's entitled docs. It trims by the real "
              "signed-in user, so the persona dropdown doesn't change it. GA trimming builds "
              "the filter from the selected persona's group, giving per-persona subsets for one "
              "tester. Verified: the compliance full-control user sees all 8 docs including "
              "MNPI. Docs must carry real tenant group/user IDs.")


# ============================================================================
# SLIDE 16 — Availability GA vs Preview
# ============================================================================
def s_availability():
    s = slide()
    bg(s, LIGHT)
    header(s, "Availability", "GA vs Preview at a glance (verified 2026-07)")
    table(s, 0.7, 1.7, 11.9, 4.2,
          ["Capability", "Status"],
          [["Microsoft Foundry portal (new)", "GA"],
           ["Foundry Hosted Agents (managed hosting)", "GA"],
           ["Publish Foundry agent \u2192 Teams & M365 Copilot", "GA"],
           ["Responses \u2192 Activity protocol bridge", "GA (automatic)"],
           ["M365 Agents Toolkit / SDK proxy path", "GA tooling"],
           ["Copilot Studio \u2192 Teams channel", "GA"],
           ["A2A (agent-to-agent) endpoint protocol", "Preview"],
           ["MCP endpoint protocol on the agent", "Preview"],
           ["Python agent-framework-foundry-hosting package", "Preview / prerelease"],
           ["M365 Agents Toolkit publishing in GCC / Gov", "Not supported"]],
          col_ratios=[3.6, 1.1], header_fill=PRIMARY, font_size=12,
          cell_colors={(6, 1): RGBColor(0xFB, 0xEE, 0xD5), (7, 1): RGBColor(0xFB, 0xEE, 0xD5),
                       (8, 1): RGBColor(0xFB, 0xEE, 0xD5), (9, 1): RGBColor(0xF7, 0xDE, 0xDE),
                       (0, 1): RGBColor(0xDE, 0xF2, 0xE6), (1, 1): RGBColor(0xDE, 0xF2, 0xE6),
                       (2, 1): RGBColor(0xDE, 0xF2, 0xE6), (3, 1): RGBColor(0xDE, 0xF2, 0xE6),
                       (4, 1): RGBColor(0xDE, 0xF2, 0xE6), (5, 1): RGBColor(0xDE, 0xF2, 0xE6)})
    text(s, 0.7, 6.05, 11.9, 0.7,
         [{"segs": [{"text": "Bottom line:  ", "bold": True, "size": 13.5, "color": OPT_B},
                    {"text": "the primary path \u2014 publish a Foundry agent directly to Teams "
                     "\u2014 is GA. Only the newer protocol surfaces (A2A, MCP) and some Python "
                     "hosting glue are Preview.", "size": 13.5, "color": DARKTEXT}]}])
    footer(s, 16)
    notes(s, "Availability snapshot as of 2026-07. The core publish experience and both primary "
              "paths are GA. A2A and MCP endpoint protocols are Preview, the Python hosting "
              "integration package is prerelease, and Agents Toolkit publishing is unsupported "
              "in Government tenants. Always re-verify before a production go/no-go.")


# ============================================================================
# SLIDE 17 — Demo runbook
# ============================================================================
def s_runbook():
    s = slide()
    bg(s, LIGHT)
    header(s, "Showcase runbook", "Live demo script \u2014 what to click, what to say",
           accent=OPT_B)
    steps = [
        ("0  Offline start", "run_all.bat \u2192 backend :8010, frontend :5174 (no Azure needed)"),
        ("1  Pick a persona", "Equity Research / Fixed-Income PM / Compliance Officer"),
        ("2  Run both options", "Ask one question \u2192 compare view + document-access map"),
        ("3  Show the trim", "Per doc: what A exposed vs what B correctly trimmed for that user"),
        ("4  Teams Option A", "Published agent answers \u201cno entitled research\u201d \u2014 explain why", AMBER),
        ("5  Teams Option B", "Sign in \u2192 SSO + OBO \u2192 entitled docs + Adaptive Card", OPT_B),
    ]
    flow(s, 0.7, 1.7, 7.5, steps, OPT_B, step_h=0.68, gap=0.13)
    card_box(s, 8.55, 1.7, 4.05, 4.55, "Talk track & fallbacks", [
        {"text": "The teaching point", "size": 12.5, "bold": True, "color": OPT_B},
        {"text": "A hosted agent AND an in-agent OBO", "size": 11},
        {"text": "tool are mutually exclusive today \u2014", "size": 11},
        {"text": "the gateway strips the user token.", "size": 11},
        {"text": " ", "size": 5},
        {"text": "Option B closes the gap", "size": 12, "bold": True, "color": PRIMARY},
        {"text": "by moving trimming to the proxy", "size": 11},
        {"text": "trust boundary.", "size": 11},
        {"text": " ", "size": 5},
        {"text": "Fallbacks", "size": 12, "bold": True, "color": AMBER},
        {"text": "\u2022 Cold agent \u2192 offline synth answer", "size": 10.5},
        {"text": "\u2022 No SSO \u2192 public-only + sign-in hint", "size": 10.5},
        {"text": "\u2022 Compliance user sees all 8 docs", "size": 10.5},
    ], PRIMARY)
    footer(s, 17)
    notes(s, "Suggested live flow. Start fully offline with run_all.bat. Pick a persona, run "
              "both options at once, and use the compare view + document-access map to show "
              "per-doc what A exposed vs what B trimmed. Then switch to Teams: Option A returns "
              "'no entitled research' (explain the token-strip), Option B signs in and returns "
              "entitled docs with an Adaptive Card. Have the offline-synth and fail-closed "
              "fallbacks ready.")


# ============================================================================
# SLIDE 18 — Takeaways
# ============================================================================
def s_takeaways():
    s = slide()
    bg(s, NAVY)
    rect(s, 0, 0, 13.333, 1.28, fill=NAVY)
    rect(s, 0.55, 0.30, 0.14, 0.66, fill=OPT_B)
    text(s, 0.85, 0.28, 11.8, 0.34, "KEY TAKEAWAYS", size=13, color=OPT_B)
    text(s, 0.85, 0.6, 11.8, 0.6, [{"text": "One agent, two front doors, one decision",
         "bold": True, "size": 26, "color": WHITE}])
    cards = [
        ("Default to Option A", "GA, lowest effort, Foundry-native tracing and versioning. "
         "Right for internal Q&A / research where users share entitlements.", OPT_A),
        ("Reserve Option B", "For per-user document security, custom SSO/OBO, and rich Teams UX "
         "(cards, buttons, HITL). The proxy is the trust boundary.", OPT_B),
        ("Per-user OBO needs a proxy", "A hosted agent's gateway strips the user token, so "
         "in-agent OBO is impossible. Trimming must move to proxy or backend.", AMBER),
        ("Both surfaces converge", "Teams and Web both resolve to the same per_user_search "
         "call. Native ACL trims by the real user; GA trimming by persona.", PRIMARY),
    ]
    x0, y0, cw, ch = 0.7, 1.65, 5.85, 2.25
    for i, (t, b, col) in enumerate(cards):
        cx = x0 + (i % 2) * (cw + 0.35)
        cy = y0 + (i // 2) * (ch + 0.3)
        rect(s, cx, cy, cw, ch, fill=NAVY2, line=col, line_w=1.5,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
        rect(s, cx, cy, 0.12, ch, fill=col)
        text(s, cx + 0.35, cy + 0.22, cw - 0.6, 0.5,
             [{"text": t, "bold": True, "size": 17, "color": col}])
        text(s, cx + 0.35, cy + 0.82, cw - 0.65, ch - 1.0,
             [{"text": b, "size": 13, "color": RGBColor(0xD5, 0xDF, 0xEA)}])
    text(s, 0.7, 6.75, 11.9, 0.5,
         [{"segs": [{"text": "Recommendation:  ", "bold": True, "size": 14, "color": OPT_B},
                    {"text": "default Option A for the fleet; hybrid-in Option B for the subset "
                     "that needs per-user security or rich UX.", "size": 14,
                     "color": RGBColor(0xBF, 0xD3, 0xE8)}]}])
    notes(s, "Four takeaways and the recommendation: default the fleet to Option A, and "
              "hybrid-in Option B only for agents that genuinely need per-user document "
              "security, custom SSO/OBO, or rich authored UX.")


# ============================================================================
# SLIDE 19 — Appendix: repo map + gotchas
# ============================================================================
def s_appendix():
    s = slide()
    bg(s, LIGHT)
    header(s, "Appendix", "Repo map & gotchas cheat-sheet")
    table(s, 0.7, 1.7, 6.7, 3.7,
          ["Component", "Path / resource"],
          [["Teams proxy (CEA)", "proxy/src/agent.ts"],
           ["Backend invocation", "backend/app/services/invoke_service.py"],
           ["Per-request identity", "backend/app/services/auth_context.py"],
           ["OBO exchange (web)", "backend/app/services/obo_service.py"],
           ["Search (app-only + per-user)", "backend/app/services/search_service.py"],
           ["Hosted agent container", "agent/app.py + agent/tools/search_tool.py"],
           ["Web SPA + MSAL", "frontend/src/auth/msalConfig.ts"],
           ["Azure Bot + OAuth", "capmarkets-obo-bot / search-sso"]],
          col_ratios=[1.5, 2.6], header_fill=PRIMARY, font_size=11)
    card_box(s, 7.65, 1.7, 4.95, 3.7, "Gotchas to remember", [
        {"text": "\u2022 Backend on port 8010 (not 8000).", "size": 11.5},
        {"text": "\u2022 Standalone bot SSO needs", "size": 11.5},
        {"text": "   api://botid-{clientId} everywhere.", "size": 11, "font": MONO},
        {"text": "\u2022 SSO connection scope = app's own", "size": 11.5},
        {"text": "   identity, OBO scope on exchangeToken.", "size": 11},
        {"text": "\u2022 Model must support encrypted content", "size": 11.5},
        {"text": "   (GPT-5 / reasoning) \u2014 gpt-4o 400s.", "size": 11},
        {"text": "\u2022 Hosted agent protocol version 2.0.0.", "size": 11.5},
        {"text": "\u2022 Portal Status N/A = scale-to-zero,", "size": 11.5},
        {"text": "   test by invoke, not portal status.", "size": 11},
    ], AMBER)
    text(s, 0.7, 5.65, 11.9, 1.0,
         [{"segs": [{"text": "Reference:  ", "bold": True, "size": 12.5, "color": PRIMARY},
                    {"text": "docs/end-to-end-flows.md (full sequence diagrams) \u00b7 "
                     "foundry-hosted-agents-teams-research.md (options A\u2013D, citations) \u00b7 "
                     "docs/teams-sso-obo-troubleshooting.md", "size": 12.5, "color": DARKTEXT}]}])
    footer(s, 19)
    notes(s, "Quick reference: where each piece lives in the repo and the hard-won gotchas "
              "(port 8010, botid- SSO form, SSO/OBO scope split, model must support encrypted "
              "content, protocol 2.0.0, scale-to-zero status). Full detail in the linked docs.")


# ----------------------------------------------------------------------------
# Build
# ----------------------------------------------------------------------------
def main():
    s_title()
    s_agenda()
    s_bigpicture()
    s_optA_overview()
    s_optA_teams()
    s_optA_web()
    s_optA_tech()
    s_optB_overview()
    s_optB_teams()
    s_optB_web()
    s_converge()
    s_optB_tech()
    s_matrix()
    s_ux()
    s_security()
    s_availability()
    s_runbook()
    s_takeaways()
    s_appendix()

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "TeamsAgents-Technical-Showcase.pptx")
    prs.save(out)
    print("Saved:", out)
    print("Slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
